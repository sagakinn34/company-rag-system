"""
Google Drive データ処理モジュール
Google Drive から各種文書を取得し、テキストを抽出します
PDF、Word、Excel、PowerPoint に対応
"""

import os
import json
import io
from typing import List, Dict, Optional
from googleapiclient.discovery import build
from google.oauth2.service_account import Credentials

# PDF/Office文書処理用ライブラリ
import PyPDF2
from docx import Document
import openpyxl
from pptx import Presentation

class GoogleDriveProcessor:
    def __init__(self, credentials_path: str):
        """
        Google Drive プロセッサーを初期化
        
        Args:
            credentials_path (str): Google Drive API認証ファイルのパス
        """
        self.credentials_path = credentials_path
        self.service = None
        self.setup_service()
    
    def setup_service(self):
        """Google Drive API サービスを設定"""
        try:
            # 必要なスコープ
            SCOPES = [
                'https://www.googleapis.com/auth/drive.readonly',
                'https://www.googleapis.com/auth/documents.readonly',
                'https://www.googleapis.com/auth/spreadsheets.readonly',
                'https://www.googleapis.com/auth/presentations.readonly'
            ]
            
            # 認証情報を読み込み
            credentials = Credentials.from_service_account_file(
                self.credentials_path, 
                scopes=SCOPES
            )
            
            # Drive API サービスを構築
            self.service = build('drive', 'v3', credentials=credentials)
            print("✅ Google Drive API サービスが正常に初期化されました")
            
        except Exception as e:
            print(f"❌ Google Drive API 初期化エラー: {e}")
            self.service = None
    
    def list_files(self, folder_id: Optional[str] = None, file_types: Optional[List[str]] = None) -> List[Dict]:
        """
        Google Drive のファイル一覧を取得
        
        Args:
            folder_id (str, optional): 特定フォルダのID
            file_types (List[str], optional): 取得するファイルタイプ
            
        Returns:
            List[Dict]: ファイル情報のリスト
        """
        if not self.service:
            return []
        
        try:
            # クエリを構築
            query_parts = []
            
            if folder_id:
                query_parts.append(f"'{folder_id}' in parents")
            
            if file_types:
                mime_types = []
                for file_type in file_types:
                    if file_type == 'pdf':
                        mime_types.append("mimeType='application/pdf'")
                    elif file_type == 'docx':
                        mime_types.append("mimeType='application/vnd.openxmlformats-officedocument.wordprocessingml.document'")
                    elif file_type == 'xlsx':
                        mime_types.append("mimeType='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'")
                    elif file_type == 'pptx':
                        mime_types.append("mimeType='application/vnd.openxmlformats-officedocument.presentationml.presentation'")
                    elif file_type == 'docs':
                        mime_types.append("mimeType='application/vnd.google-apps.document'")
                    elif file_type == 'sheets':
                        mime_types.append("mimeType='application/vnd.google-apps.spreadsheet'")
                    elif file_type == 'slides':
                        mime_types.append("mimeType='application/vnd.google-apps.presentation'")
                
                if mime_types:
                    query_parts.append(f"({' or '.join(mime_types)})")
            
            query_parts.append("trashed=false")
            query = ' and '.join(query_parts)
            
            # ファイル一覧を取得
            results = self.service.files().list(
                q=query,
                pageSize=100,
                fields="nextPageToken, files(id, name, mimeType, size, createdTime, modifiedTime)"
            ).execute()
            
            files = results.get('files', [])
            print(f"📁 {len(files)} 件のファイルが見つかりました")
            
            return files
            
        except Exception as e:
            print(f"❌ ファイル一覧取得エラー: {e}")
            return []
    
    def extract_text_from_pdf(self, file_content: bytes) -> str:
        """PDFファイルからテキストを抽出"""
        try:
            pdf_file = io.BytesIO(file_content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
            
            return text.strip()
        except Exception as e:
            print(f"❌ PDF テキスト抽出エラー: {e}")
            return ""
    
    def extract_text_from_docx(self, file_content: bytes) -> str:
        """Word文書からテキストを抽出"""
        try:
            doc_file = io.BytesIO(file_content)
            doc = Document(doc_file)
            
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            
            return text.strip()
        except Exception as e:
            print(f"❌ Word テキスト抽出エラー: {e}")
            return ""
    
    def extract_text_from_xlsx(self, file_content: bytes) -> str:
        """Excelファイルからテキストを抽出"""
        try:
            excel_file = io.BytesIO(file_content)
            workbook = openpyxl.load_workbook(excel_file)
            
            text = ""
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"=== {sheet_name} ===\n"
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = []
                    for cell in row:
                        if cell is not None:
                            row_text.append(str(cell))
                    if row_text:
                        text += "\t".join(row_text) + "\n"
                
                text += "\n"
            
            return text.strip()
        except Exception as e:
            print(f"❌ Excel テキスト抽出エラー: {e}")
            return ""
    
    def extract_text_from_pptx(self, file_content: bytes) -> str:
        """PowerPointファイルからテキストを抽出"""
        try:
            ppt_file = io.BytesIO(file_content)
            presentation = Presentation(ppt_file)
            
            text = ""
            for i, slide in enumerate(presentation.slides):
                text += f"=== スライド {i+1} ===\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                
                text += "\n"
            
            return text.strip()
        except Exception as e:
            print(f"❌ PowerPoint テキスト抽出エラー: {e}")
            return ""
    
    def download_and_extract_text(self, file_id: str, mime_type: str) -> str:
        """ファイルをダウンロードしてテキストを抽出"""
        if not self.service:
            return ""
        
        try:
            # Google Docs形式の場合はエクスポート
            if mime_type == 'application/vnd.google-apps.document':
                request = self.service.files().export_media(fileId=file_id, mimeType='text/plain')
            elif mime_type == 'application/vnd.google-apps.spreadsheet':
                request = self.service.files().export_media(fileId=file_id, mimeType='text/csv')
            elif mime_type == 'application/vnd.google-apps.presentation':
                request = self.service.files().export_media(fileId=file_id, mimeType='text/plain')
            else:
                # その他のファイルは直接ダウンロード
                request = self.service.files().get_media(fileId=file_id)
            
            file_content = request.execute()
            
            # ファイルタイプに応じてテキスト抽出
            if mime_type == 'application/pdf':
                return self.extract_text_from_pdf(file_content)
            elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
                return self.extract_text_from_docx(file_content)
            elif mime_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
                return self.extract_text_from_xlsx(file_content)
            elif mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
                return self.extract_text_from_pptx(file_content)
            elif 'google-apps' in mime_type:
                # Google Docs形式はプレーンテキストとして処理
                return file_content.decode('utf-8')
            else:
                # その他はプレーンテキストとして試行
                try:
                    return file_content.decode('utf-8')
                except:
                    return file_content.decode('utf-8', errors='ignore')
        
        except Exception as e:
            print(f"❌ ファイルダウンロード・テキスト抽出エラー: {e}")
            return ""
    
    def process_all_files(self) -> List[Dict]:
        """
        全ファイルを処理してテキストを抽出
        
        Returns:
            List[Dict]: 処理されたドキュメントのリスト
        """
        print("🔍 Google Drive ファイル処理を開始...")
        
        # 対応ファイルタイプ
        supported_types = ['pdf', 'docx', 'xlsx', 'pptx', 'docs', 'sheets', 'slides']
        files = self.list_files(file_types=supported_types)
        
        documents = []
        
        for file_info in files:
            print(f"📄 処理中: {file_info['name']}")
            
            # テキストを抽出
            text = self.download_and_extract_text(file_info['id'], file_info['mimeType'])
            
            if text.strip():
                document = {
                    'id': f"gdrive_{file_info['id']}",
                    'title': file_info['name'],
                    'content': text,
                    'source': 'google_drive',
                    'mime_type': file_info['mimeType'],
                    'size': file_info.get('size', '0'),
                    'created_time': file_info.get('createdTime', ''),
                    'modified_time': file_info.get('modifiedTime', ''),
                    'url': f"https://drive.google.com/file/d/{file_info['id']}/view"
                }
                documents.append(document)
                print(f"✅ テキスト抽出成功: {len(text)} 文字")
            else:
                print(f"⚠️ テキスト抽出失敗: {file_info['name']}")
        
        print(f"🎉 Google Drive 処理完了: {len(documents)} 件のドキュメント")
        return documents

def test_google_drive_processor():
    """Google Drive プロセッサーのテスト実行"""
    print("=== Google Drive プロセッサー テスト実行 ===")
    
    # 認証ファイルのパス
    credentials_path = "./credentials/google-drive-credentials.json"
    
    if not os.path.exists(credentials_path):
        print(f"❌ 認証ファイルが見つかりません: {credentials_path}")
        return
    
    try:
        # プロセッサーを初期化
        processor = GoogleDriveProcessor(credentials_path)
        
        # ファイル一覧を取得（テスト）
        files = processor.list_files()
        print(f"📁 利用可能ファイル数: {len(files)}")
        
        if files:
            print("\n📋 ファイル一覧（最初の5件）:")
            for i, file_info in enumerate(files[:5]):
                print(f"  {i+1}. {file_info['name']} ({file_info['mimeType']})")
        
        # 実際の処理はコメントアウト（テスト時に大量データ処理を避けるため）
        # documents = processor.process_all_files()
        
        print("\n✅ Google Drive プロセッサーテスト完了")
        
    except Exception as e:
        print(f"❌ テスト実行エラー: {e}")

if __name__ == "__main__":
    test_google_drive_processor()
